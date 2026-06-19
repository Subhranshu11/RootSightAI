import os
import glob
import json
import pickle
import pandas as pd
import faiss

from sentence_transformers import SentenceTransformer

from docx import Document
from pptx import Presentation
from PyPDF2 import PdfReader
import xml.etree.ElementTree as ET

# --------------------------------------------------
# CONFIG
# --------------------------------------------------

DATA_FOLDER = "data"
VECTORSTORE_FOLDER = "vectorstore"

FAISS_INDEX_PATH = os.path.join(
    VECTORSTORE_FOLDER,
    "faiss_index.bin"
)

METADATA_PATH = os.path.join(
    VECTORSTORE_FOLDER,
    "metadata.pkl"
)

EMBEDDING_MODEL = "all-MiniLM-L6-v2"

CHUNK_SIZE = 1200
OVERLAP = 200

# --------------------------------------------------
# LOAD MODEL
# --------------------------------------------------

print("Loading embedding model...")

model = SentenceTransformer(
    EMBEDDING_MODEL
)

# --------------------------------------------------
# SMART CHUNKING
# --------------------------------------------------

def chunk_text(text):

    text = str(text).strip()

    if not text:
        return []

    chunks = []

    start = 0

    while start < len(text):

        end = start + CHUNK_SIZE

        chunks.append(
            text[start:end]
        )

        start += (
            CHUNK_SIZE - OVERLAP
        )

    return chunks

# --------------------------------------------------
# FILE READERS
# --------------------------------------------------

def read_csv(path):

    chunks = []

    df = pd.read_csv(path)

    for _, row in df.iterrows():

        row_text = "\n".join(
            [
                f"{col}: {row[col]}"
                for col in df.columns
            ]
        )

        chunks.extend(
            chunk_text(row_text)
        )

    return chunks


def read_xlsx(path):

    chunks = []

    sheets = pd.read_excel(
        path,
        sheet_name=None
    )

    for sheet_name, df in sheets.items():

        for _, row in df.iterrows():

            row_text = "\n".join(
                [
                    f"{col}: {row[col]}"
                    for col in df.columns
                ]
            )

            chunks.extend(
                chunk_text(row_text)
            )

    return chunks


def read_json(path):

    with open(
        path,
        "r",
        encoding="utf-8"
    ) as f:

        data = json.load(f)

    text = json.dumps(
        data,
        indent=2
    )

    return chunk_text(text)


def read_txt(path):

    with open(
        path,
        "r",
        encoding="utf-8",
        errors="ignore"
    ) as f:

        return chunk_text(
            f.read()
        )


def read_md(path):

    return read_txt(path)


def read_log(path):

    return read_txt(path)


def read_xml(path):

    tree = ET.parse(path)

    root = tree.getroot()

    text = ET.tostring(
        root,
        encoding="unicode"
    )

    return chunk_text(text)


def read_docx(path):

    doc = Document(path)

    text = "\n".join(
        [
            para.text
            for para in doc.paragraphs
        ]
    )

    return chunk_text(text)


def read_pdf(path):

    reader = PdfReader(path)

    text = ""

    for page in reader.pages:

        try:
            text += (
                page.extract_text()
                + "\n"
            )

        except:
            pass

    return chunk_text(text)


def read_pptx(path):

    prs = Presentation(path)

    text = ""

    for slide in prs.slides:

        for shape in slide.shapes:

            if hasattr(
                shape,
                "text"
            ):

                text += (
                    shape.text
                    + "\n"
                )

    return chunk_text(text)

# --------------------------------------------------
# INGEST ALL FILES
# --------------------------------------------------

all_chunks = []

SUPPORTED_PATTERNS = [
    "*.csv",
    "*.xlsx",
    "*.json",
    "*.txt",
    "*.md",
    "*.log",
    "*.xml",
    "*.docx",
    "*.pdf",
    "*.pptx"
]

all_files = []

for pattern in SUPPORTED_PATTERNS:

    all_files.extend(
        glob.glob(
            os.path.join(
                DATA_FOLDER,
                pattern
            )
        )
    )

print(
    f"Found {len(all_files)} files"
)

for file_path in all_files:

    try:

        ext = (
            os.path.splitext(
                file_path
            )[1]
            .lower()
        )

        print(
            f"\nProcessing: {file_path}"
        )

        chunks = []

        if ext == ".csv":
            chunks = read_csv(file_path)

        elif ext == ".xlsx":
            chunks = read_xlsx(file_path)

        elif ext == ".json":
            chunks = read_json(file_path)

        elif ext == ".txt":
            chunks = read_txt(file_path)

        elif ext == ".md":
            chunks = read_md(file_path)

        elif ext == ".log":
            chunks = read_log(file_path)

        elif ext == ".xml":
            chunks = read_xml(file_path)

        elif ext == ".docx":
            chunks = read_docx(file_path)

        elif ext == ".pdf":
            chunks = read_pdf(file_path)

        elif ext == ".pptx":
            chunks = read_pptx(file_path)

        print(
            f"Chunks Created: {len(chunks)}"
        )

        all_chunks.extend(
            chunks
        )

    except Exception as e:

        print(
            f"Error processing {file_path}"
        )

        print(e)

# --------------------------------------------------
# CREATE EMBEDDINGS
# --------------------------------------------------

print(
    f"\nCreating embeddings for {len(all_chunks)} chunks..."
)

embeddings = model.encode(
    all_chunks,
    show_progress_bar=True
)

embeddings = embeddings.astype(
    "float32"
)

# --------------------------------------------------
# CREATE FAISS INDEX
# --------------------------------------------------

dimension = embeddings.shape[1]

index = faiss.IndexFlatL2(
    dimension
)

index.add(
    embeddings
)

# --------------------------------------------------
# SAVE
# --------------------------------------------------

os.makedirs(
    VECTORSTORE_FOLDER,
    exist_ok=True
)

faiss.write_index(
    index,
    FAISS_INDEX_PATH
)

with open(
    METADATA_PATH,
    "wb"
) as f:

    pickle.dump(
        all_chunks,
        f
    )

# --------------------------------------------------
# DONE
# --------------------------------------------------

print("\n=================================")
print("VECTOR DATABASE CREATED")
print(f"Chunks: {len(all_chunks)}")
print(f"Index: {FAISS_INDEX_PATH}")
print(f"Metadata: {METADATA_PATH}")
print("=================================")